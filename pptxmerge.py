import copy
import io
import re
from typing import List, Dict, Optional
from PIL import Image
from pptx import Presentation

TAG_RE = re.compile(r'<<\s*([A-Z0-9_]+)\s*>>')


def _duplicate_slide(prs: Presentation, source_slide):
    layout = prs.slide_layouts[0] if len(prs.slide_layouts) else None
    new_slide = prs.slides.add_slide(layout) if layout is not None else prs.slides.add_slide(prs.slide_layouts[len(prs.slide_layouts)-1])
    for shp in source_slide.shapes:
        newel = copy.deepcopy(shp.element)
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    return new_slide


def _shape_text(shape):
    if not getattr(shape, 'has_text_frame', False):
        return ''
    return '\n'.join(''.join(r.text for r in p.runs) if p.runs else p.text for p in shape.text_frame.paragraphs)


def _replace_text_in_shape(shape, values: Dict[str, str]):
    if not getattr(shape, 'has_text_frame', False):
        return False
    changed = False
    tf = shape.text_frame
    for para in tf.paragraphs:
        runs = para.runs
        if not runs:
            continue
        full_text = ''.join(r.text for r in runs)
        if '<<' not in full_text:
            continue

        def repl(m):
            key = m.group(1)
            return str(values.get(key, m.group(0)) or '')

        new_text = TAG_RE.sub(repl, full_text)
        if new_text != full_text:
            runs[0].text = new_text
            for run in runs[1:]:
                run.text = ''
            changed = True
    return changed


def _find_tag_shape(slide, tag_name: str):
    pat = re.compile(r'<<\s*' + re.escape(tag_name) + r'\s*>>')
    for shp in slide.shapes:
        if getattr(shp, 'has_text_frame', False) and pat.search(_shape_text(shp)):
            return shp
    return None


def _crop_cover(img: Image.Image, width: int, height: int):
    img = img.convert('RGB')
    tr = width / height
    ir = img.width / img.height if img.height else tr
    if ir > tr:
        nw = int(img.height * tr)
        x0 = max((img.width - nw) // 2, 0)
        img = img.crop((x0, 0, x0 + nw, img.height))
    else:
        nh = int(img.width / tr) if tr else img.height
        y0 = max((img.height - nh) // 2, 0)
        img = img.crop((0, y0, img.width, y0 + nh))
    return img


def _replace_photo_shape(slide, shp, pil_image: Image.Image):
    left, top, width, height = shp.left, shp.top, shp.width, shp.height
    img = _crop_cover(pil_image, width, height)
    buf = io.BytesIO()
    img.save(buf, format='PNG')
    buf.seek(0)
    slide.shapes.add_picture(buf, left, top, width=width, height=height)
    try:
        shp._element.getparent().remove(shp._element)
    except Exception:
        pass


def build_merged_presentation(template_pptx_path_or_file, slide_groups: List[Dict], processed_photo_lookup: Optional[Dict[str, Image.Image]] = None):
    processed_photo_lookup = processed_photo_lookup or {}
    prs = Presentation(template_pptx_path_or_file)
    if not prs.slides:
        raise ValueError('Template has no slides.')
    template_slide = prs.slides[0]
    for group in slide_groups:
        slide = _duplicate_slide(prs, template_slide)
        speakers = group.get('speakers', [])
        values = {
            'SESSION_NAME': str(group.get('SESSION_NAME', '') or ''),
            'HALL_NAME': str(group.get('HALL_NAME', '') or ''),
            'DATE': str(group.get('DATE', '') or ''),
            'MAIN_SESSION_DETAILS': str(group.get('MAIN_SESSION_DETAILS', '') or ''),
            'SPEAKER_SESSION_DETAILS': str(group.get('SPEAKER_SESSION_DETAILS', '') or ''),
            'PLACEHOLDER_1': str(group.get('PLACEHOLDER_1', '') or ''),
            'PLACEHOLDER_2': str(group.get('PLACEHOLDER_2', '') or ''),
        }
        for i, sp in enumerate(speakers, start=1):
            values[f'SPEAKER_NAME_{i}'] = str(sp.get('name', '') or '')
            values[f'SPEAKER_TITLE_{i}'] = str(sp.get('title', '') or '')
            values[f'SPEAKER_COMPANY_{i}'] = str(sp.get('company', '') or '')
        for shp in list(slide.shapes):
            _replace_text_in_shape(shp, values)
        for i, sp in enumerate(speakers, start=1):
            key = sp.get('photo_key')
            if key and key in processed_photo_lookup:
                photo_shape = _find_tag_shape(slide, f'SPEAKER_PHOTO_{i}')
                if photo_shape is not None:
                    _replace_photo_shape(slide, photo_shape, processed_photo_lookup[key])
    return prs


def save_presentation_to_bytes(prs: Presentation) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
