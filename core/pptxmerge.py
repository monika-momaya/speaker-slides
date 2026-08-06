import copy
import io
import re
from typing import List, Dict, Optional
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE

TAG_RE = re.compile(r'<<\s*([A-Z0-9_]+)\s*>>')


def _duplicate_slide(prs: Presentation, source_slide):
    layout = prs.slide_layouts[0] if len(prs.slide_layouts) else None
    new_slide = prs.slides.add_slide(layout) if layout is not None else prs.slides.add_slide(prs.slide_layouts[len(prs.slide_layouts)-1])
    for shp in source_slide.shapes:
        newel = copy.deepcopy(shp.element)
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    return new_slide


def _shape_text(shape):
    if not getattr(shape, 'has_text_frame', False):
        return ''
    return '\n'.join(''.join(r.text for r in p.runs) if p.runs else p.text for p in shape.text_frame.paragraphs)


def _replace_text_in_shape(shape, values: Dict[str, str]):
    if not getattr(shape, 'has_text_frame', False):
        return False
    changed = False
    tf = shape.text_frame
    for para in tf.paragraphs:
        runs = para.runs
        if not runs:
            continue
        full_text = ''.join(r.text for r in runs)
        if '<<' not in full_text:
            continue

        def repl(m):
            key = m.group(1)
            return str(values.get(key, m.group(0)) or '')

        new_text = TAG_RE.sub(repl, full_text)
        if new_text != full_text:
            runs[0].text = new_text
            for run in runs[1:]:
                run.text = ''
            changed = True
    return changed


def _find_tag_shape(slide, tag_name: str):
    pat = re.compile(r'<<\s*' + re.escape(tag_name) + r'\s*>>')
    for shp in slide.shapes:
        if getattr(shp, 'has_text_frame', False) and pat.search(_shape_text(shp)):
            return shp
    return None


def _placeholder_mask_kind(shp):
    """Detect the placeholder shape's outline so the inserted photo can
    be masked to match, instead of always being a plain hard-edged
    rectangle that just covers the shape's bounding box."""
    try:
        if shp.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if shp.auto_shape_type == MSO_AUTO_SHAPE_TYPE.OVAL:
                return 'oval'
            if shp.auto_shape_type == MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE:
                return 'rounded_rect'
    except Exception:
        pass
    return None


def _rounded_rect_corner_radius(shp, img_w: int, img_h: int) -> int:
    """Approximate the corner radius (in output pixels) PowerPoint would
    render for this rounded-rectangle shape, using its adjustment value
    if available, else a sensible default proportion."""
    try:
        adj = shp.adjustments[0]
    except Exception:
        adj = 0.16  # PowerPoint's default rounded-rectangle adjustment
    short_side = min(img_w, img_h)
    return max(1, int(short_side * adj))


def _mask_to_oval(img: Image.Image, supersample: int = 4) -> Image.Image:
    """Return an RGBA copy of img with an elliptical alpha mask applied
    (inscribed in img's full bounding box), antialiased via supersampling."""
    w, h = img.size
    big_w, big_h = w * supersample, h * supersample
    mask = Image.new('L', (big_w, big_h), 0)
    draw = ImageDraw.Draw(mask)
    draw.ellipse((0, 0, big_w - 1, big_h - 1), fill=255)
    mask = mask.resize((w, h), Image.LANCZOS)
    img = img.convert('RGBA')
    img.putalpha(mask)
    return img


def _mask_to_rounded_rect(img: Image.Image, radius: int, supersample: int = 4) -> Image.Image:
    """Return an RGBA copy of img with a rounded-rectangle alpha mask
    applied, antialiased via supersampling."""
    w, h = img.size
    big_w, big_h = w * supersample, h * supersample
    mask = Image.new('L', (big_w, big_h), 0)
    draw = ImageDraw.Draw(mask)
    draw.rounded_rectangle((0, 0, big_w - 1, big_h - 1), radius=radius * supersample, fill=255)
    mask = mask.resize((w, h), Image.LANCZOS)
    img = img.convert('RGBA')
    img.putalpha(mask)
    return img


def _crop_cover(img: Image.Image, width: int, height: int):
    img = img.convert('RGB')
    tr = width / height
    ir = img.width / img.height if img.height else tr
    if ir > tr:
        nw = int(img.height * tr)
        x0 = max((img.width - nw) // 2, 0)
        img = img.crop((x0, 0, x0 + nw, img.height))
    else:
        nh = int(img.width / tr) if tr else img.height
        y0 = max((img.height - nh) // 2, 0)
        img = img.crop((0, y0, img.width, y0 + nh))
    return img


def _replace_photo_shape(slide, shp, pil_image: Image.Image):
    left, top, width, height = shp.left, shp.top, shp.width, shp.height
    mask_kind = _placeholder_mask_kind(shp)
    img = _crop_cover(pil_image, width, height)
    if mask_kind == 'oval':
        img = _mask_to_oval(img)
    elif mask_kind == 'rounded_rect':
        # width/height are EMUs; use the pixel size of img (already
        # matches the box's aspect ratio) to compute a proportional radius
        radius = _rounded_rect_corner_radius(shp, img.width, img.height)
        img = _mask_to_rounded_rect(img, radius)
    buf = io.BytesIO()
    img.save(buf, format='PNG')
    buf.seek(0)
    slide.shapes.add_picture(buf, left, top, width=width, height=height)
    try:
        shp._element.getparent().remove(shp._element)
    except Exception:
        pass


LEFTOVER_TEXT_TAG_RE = re.compile(r'<<\s*SPEAKER_(?:NAME|TITLE|COMPANY)_(\d+)\s*>>')
LEFTOVER_PHOTO_TAG_RE = re.compile(r'<<\s*SPEAKER_PHOTO_(\d+)\s*>>')


def _clean_unused_speaker_slots(slide):
    """After the known speakers for this group have been filled in, some
    slots on a multi-speaker template may be left over (e.g. a 4-speaker
    template used for a 2-speaker slide). Blank any leftover
    SPEAKER_NAME/TITLE/COMPANY_n text and remove any leftover
    SPEAKER_PHOTO_n placeholder shape entirely, so unused slots disappear
    cleanly instead of showing raw tag text or an empty circle."""
    for shp in list(slide.shapes):
        if not getattr(shp, 'has_text_frame', False):
            continue
        text = _shape_text(shp)
        if LEFTOVER_PHOTO_TAG_RE.search(text):
            try:
                shp._element.getparent().remove(shp._element)
            except Exception:
                pass
            continue
        if LEFTOVER_TEXT_TAG_RE.search(text):
            for para in shp.text_frame.paragraphs:
                for run in para.runs:
                    run.text = LEFTOVER_TEXT_TAG_RE.sub('', run.text)


SPEAKER_SLOT_TAG_RE = re.compile(r'SPEAKER_(?:NAME|TITLE|COMPANY|PHOTO)_(\d+)')


def _detect_max_speaker_slot(slide) -> int:
    """Scan a template slide's shapes for the highest SPEAKER_..._n tag
    suffix present, so we know how many speakers that particular layout
    (single vs panel) can actually hold."""
    max_slot = 0
    for shp in slide.shapes:
        text = _shape_text(shp)
        for m in SPEAKER_SLOT_TAG_RE.finditer(text):
            max_slot = max(max_slot, int(m.group(1)))
    return max_slot


def _delete_slide(prs: Presentation, index: int):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    slide_id_elem = slides[index]
    rId = slide_id_elem.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    xml_slides.remove(slide_id_elem)
    try:
        prs.part.drop_rel(rId)
    except Exception:
        pass


def _base_values(group: Dict) -> Dict[str, str]:
    return {
        'SESSION_NAME': str(group.get('SESSION_NAME', '') or ''),
        'HALL_NAME': str(group.get('HALL_NAME', '') or ''),
        'DATE': str(group.get('DATE', '') or ''),
        'MAIN_SESSION_DETAILS': str(group.get('MAIN_SESSION_DETAILS', '') or ''),
        'SPEAKER_SESSION_DETAILS': str(group.get('SPEAKER_SESSION_DETAILS', '') or ''),
        'PLACEHOLDER_1': str(group.get('PLACEHOLDER_1', '') or ''),
        'PLACEHOLDER_2': str(group.get('PLACEHOLDER_2', '') or ''),
    }


def _fill_slide(slide, values: Dict[str, str], speakers: List[Dict], processed_photo_lookup: Dict[str, Image.Image]):
    """Fill a duplicated slide's text and photo tags for the given
    speakers (positioned starting at slot 1), then clean up any leftover
    unused speaker slots the layout has room for but weren't used."""
    for i, sp in enumerate(speakers, start=1):
        values[f'SPEAKER_NAME_{i}'] = str(sp.get('name', '') or '')
        values[f'SPEAKER_TITLE_{i}'] = str(sp.get('title', '') or '')
        values[f'SPEAKER_COMPANY_{i}'] = str(sp.get('company', '') or '')
    for shp in list(slide.shapes):
        _replace_text_in_shape(shp, values)
    for i, sp in enumerate(speakers, start=1):
        key = sp.get('photo_key')
        if key and key in processed_photo_lookup:
            photo_shape = _find_tag_shape(slide, f'SPEAKER_PHOTO_{i}')
            if photo_shape is not None:
                _replace_photo_shape(slide, photo_shape, processed_photo_lookup[key])
    _clean_unused_speaker_slots(slide)


def build_merged_presentation(template_pptx_path_or_file, slide_groups: List[Dict], processed_photo_lookup: Optional[Dict[str, Image.Image]] = None):
    """Generate the merged deck.

    The template pptx's slide 1 (index 0) is always the single-speaker
    layout, used for any group with exactly one speaker, and also used
    to generate one dedicated slide per speaker for multi-speaker
    (panel) groups.

    Any additional slides (index 1, 2, 3, ...) are treated as panel
    "tiers" of increasing capacity, in any order — e.g. one slide with
    room for 3 speakers, another with room for 6, another for 9. Each
    tier's capacity is auto-detected from its own SPEAKER_..._n tags.
    For a multi-speaker group, the SMALLEST tier that can still fit all
    of that group's speakers is used for the combined panel slide (so a
    4-speaker panel uses a 6-speaker tier rather than a 9-speaker one,
    if both exist, keeping the layout as tight as possible). If no tier
    is big enough, the largest available tier is used, filled to its
    capacity. Either way, every speaker in the group still gets their
    own individual slide afterwards, built from the single-speaker
    layout, so nobody is silently dropped.

    The original template slide(s) are removed from the final output so
    the deck contains only the generated content.
    """
    processed_photo_lookup = processed_photo_lookup or {}
    prs = Presentation(template_pptx_path_or_file)
    if not prs.slides:
        raise ValueError('Template has no slides.')

    num_template_slides = len(prs.slides)
    single_template_slide = prs.slides[0]

    # Build the list of (capacity, slide) tiers from any slides beyond index 0,
    # sorted smallest-capacity first so we can pick the tightest fit.
    tiers = []
    for idx in range(1, num_template_slides):
        tier_slide = prs.slides[idx]
        capacity = _detect_max_speaker_slot(tier_slide)
        if capacity > 0:
            tiers.append((capacity, tier_slide))
    tiers.sort(key=lambda t: t[0])

    def _best_tier(speaker_count: int):
        for capacity, tier_slide in tiers:
            if capacity >= speaker_count:
                return capacity, tier_slide
        return tiers[-1] if tiers else (0, None)

    for group in slide_groups:
        speakers = group.get('speakers', [])
        base_values = _base_values(group)

        if len(speakers) > 1 and tiers:
            capacity, tier_slide = _best_tier(len(speakers))
            if tier_slide is not None:
                panel_slide = _duplicate_slide(prs, tier_slide)
                _fill_slide(panel_slide, dict(base_values), speakers[:capacity], processed_photo_lookup)

        if len(speakers) <= 1:
            slide = _duplicate_slide(prs, single_template_slide)
            _fill_slide(slide, dict(base_values), speakers, processed_photo_lookup)
        else:
            # one dedicated individual slide per speaker, right after the panel slide
            for sp in speakers:
                slide = _duplicate_slide(prs, single_template_slide)
                _fill_slide(slide, dict(base_values), [sp], processed_photo_lookup)

    # Remove the original template slide(s) so only generated content remains.
    for _ in range(num_template_slides):
        _delete_slide(prs, 0)

    return prs


def save_presentation_to_bytes(prs: Presentation) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
