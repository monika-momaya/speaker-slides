import re
from dataclasses import dataclass, field
from typing import List, Dict, Set
from pptx import Presentation

TAG_PATTERN = re.compile(r"<<\s*([A-Z0-9_]+)\s*>>")
KNOWN_SINGLE_TAGS = [
    "SESSION_NAME",
    "HALL_NAME",
    "DATE",
    "MAIN_SESSION_DETAILS",
    "SPEAKER_SESSION_DETAILS",
    "PLACEHOLDER_1",
    "PLACEHOLDER_2",
]
SPEAKER_FIELD_PREFIXES = ["SPEAKER_NAME_", "SPEAKER_TITLE_", "SPEAKER_COMPANY_", "SPEAKER_PHOTO_"]

@dataclass
class TemplateTagReport:
    all_tags: Set[str] = field(default_factory=set)
    found_single_tags: List[str] = field(default_factory=list)
    missing_single_tags: List[str] = field(default_factory=list)
    max_speaker_slot: int = 0
    speaker_fields_present: Dict[str, List[int]] = field(default_factory=dict)
    warnings: List[str] = field(default_factory=list)


def _shape_text(shape):
    if not getattr(shape, "has_text_frame", False):
        return ""
    texts = []
    for p in shape.text_frame.paragraphs:
        parts = [r.text for r in p.runs] if p.runs else [p.text]
        texts.append("".join(parts))
    return "\n".join(texts)


def extract_tags_from_pptx(pptx_path_or_file) -> TemplateTagReport:
    prs = Presentation(pptx_path_or_file)
    if len(prs.slides) == 0:
        raise ValueError("Uploaded template has no slides.")
    slide = prs.slides[0]
    all_tags: Set[str] = set()
    for shape in slide.shapes:
        full_text = _shape_text(shape)
        for m in TAG_PATTERN.finditer(full_text):
            all_tags.add(m.group(1))
    report = TemplateTagReport(all_tags=all_tags)
    for tag in KNOWN_SINGLE_TAGS:
        if tag in all_tags:
            report.found_single_tags.append(tag)
        else:
            report.missing_single_tags.append(tag)
    speaker_fields_present: Dict[str, List[int]] = {p: [] for p in SPEAKER_FIELD_PREFIXES}
    max_slot = 0
    for tag in all_tags:
        for prefix in SPEAKER_FIELD_PREFIXES:
            if tag.startswith(prefix):
                suffix = tag[len(prefix):]
                if suffix.isdigit():
                    idx = int(suffix)
                    speaker_fields_present[prefix].append(idx)
                    max_slot = max(max_slot, idx)
    for k in speaker_fields_present:
        speaker_fields_present[k].sort()
    report.speaker_fields_present = speaker_fields_present
    report.max_speaker_slot = max_slot
    if max_slot == 0:
        report.warnings.append("No <<SPEAKER_NAME_n>> style tags were found on slide 1.")
    for slot in range(1, max_slot + 1):
        if slot not in speaker_fields_present["SPEAKER_PHOTO_"]:
            report.warnings.append(f"Missing <<SPEAKER_PHOTO_{slot}>> textbox placeholder.")
    return report
